# phonics_maker/pdf_generation/pptx_generator.py

"""
Generate a PowerPoint (.pptx) presentation from processed book images.
Each slide is a full-bleed image — ideal for smartboard classroom display.
"""

from typing import List
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu
from app.core.config.logger import logger


# A5 dimensions in EMUs (English Metric Units: 1mm = 36000 EMU)
# A5 = 148mm × 210mm
A5_WIDTH_EMU = Emu(5328000)   # 148mm × 36000
A5_HEIGHT_EMU = Emu(7560000)  # 210mm × 36000


async def generate_pptx_from_images(
    image_paths: List[str],
    safe_title: str,
    output_dir: Path,
) -> str:
    """
    Generate a PowerPoint file with one full-bleed image per slide.

    Args:
        image_paths: Ordered list of local image file paths (cover first, then scenes).
        safe_title: Filesystem-safe title for the output file.
        output_dir: Directory where the PPTX will be saved.

    Returns:
        File path of the generated PPTX.
    """
    try:
        logger.info(f"Generating PPTX with {len(image_paths)} slides")

        output_dir.mkdir(exist_ok=True)
        pptx_path = output_dir / f"{safe_title}.pptx"

        prs = Presentation()
        prs.slide_width = A5_WIDTH_EMU
        prs.slide_height = A5_HEIGHT_EMU

        # Use a blank slide layout (index 6 is typically blank)
        blank_layout = prs.slide_layouts[6]

        for i, img_path in enumerate(image_paths):
            # Strip file:// prefix if present (processed images use file:// URLs)
            if img_path.startswith("file://"):
                img_path = img_path[7:]  # len("file://") == 7
            if not Path(img_path).exists():
                logger.warning(f"PPTX: Skipping missing image {img_path}")
                continue

            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(
                img_path,
                left=0,
                top=0,
                width=prs.slide_width,
                height=prs.slide_height,
            )
            logger.debug(f"PPTX: Added slide {i + 1} from {Path(img_path).name}")

        prs.save(str(pptx_path))
        logger.info(f"PPTX created successfully: {pptx_path}")
        return str(pptx_path)

    except Exception as e:
        logger.error(f"Error generating PPTX: {str(e)}")
        raise e
